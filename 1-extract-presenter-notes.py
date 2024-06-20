import collections 
import collections.abc
from pptx import Presentation
import re
import sys
import os.path

if len(sys.argv) == 2:
        input_file = sys.argv[1]
else:
        print ("usage: "+sys.argv[0]+" <pptx filename>")
        exit(1)

ppt=Presentation(input_file)
notes = []
filename=os.path.basename(input_file)
output_file=filename[:-5]+"_presenter_notes.md"

f = open(output_file,"w")
f.write("# {} presenter notes and links".format(filename))
for page, slide in enumerate(ppt.slides):
    # this is the notes that doesn't appear on the ppt slide,
    # but really the 'presenter' note. 
    text_notes = slide.notes_slide.notes_text_frame.text
    try:
        # Regular expression pattern to match HTTPS links
        pattern = r'https://\S+'
        #print("TST: Slide {}, length {}".format(int(page)+1, len(slide.shapes.placeholders)))
        text_frame_text = slide.shapes.placeholders[1].text_frame.text
    except Exception as error:
        try:
            text_frame_text = slide.shapes.placeholders[0].text_frame.text
        except:
            # slide does not contain any text, only title:
            text_frame_text=""
        
    f.write("\n\n-------\n### Slide {}: {}\n\n{}\n".format(int(page)+1, text_frame_text, text_notes))
    # Use re.findall() to extract all matches
    https_links = re.findall(pattern, text_notes)
    if len(https_links) > 0:
        f.write("\nHTTPS links found:\n")
        for link in https_links:
            cleaned_link = re.sub(r'[^/\w\d]+$', '', link)
            f.write("* {}\n".format(cleaned_link))

f.close()