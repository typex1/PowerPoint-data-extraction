from pptx import Presentation
import os.path
import sys
import re

# Package "pptx" documentation: https://python-pptx.readthedocs.io/en/latest/user/notes.html

if len(sys.argv) == 2:
        pptx_file = sys.argv[1]
else:
        print ("usage: "+sys.argv[0]+" <pptx filename>")
        exit(1)

filename = os.path.basename(pptx_file)
output_file = filename[:-5]+"_presenter_notes.md"

import re

def extract_module_number(text):
    """
    Extracts the module information from a given text.

    Args:
        text (str): The input text.

    Returns:
        str: The extracted module information, or an empty string if not found.
    """
    pattern = r"Module\s+\d+"
    match = re.search(pattern, text)
    if match:
        return match.group()
    else:
        return ""

def extract_slide_content(pptx_file):
    """
    Extracts slide content (title, text, and notes) from a PowerPoint presentation.

    Args:
        pptx_file (str): Path to the PowerPoint file.

    Returns:
        A list of dictionaries, where each dictionary represents a slide and contains the following keys:
        - 'title': The title of the slide.
        - 'text': The text content of the slide.
        - 'notes': The notes associated with the slide.
    """
    presentation = Presentation(pptx_file)
    slide_contents = []

    module_no = ''
    for slide in presentation.slides:
        try:
            slide_data = {
                'title': slide.shapes.title.text if slide.shapes.title else '',
                'text': '\n'.join(shape.text for shape in slide.shapes if hasattr(shape, 'text')),
                'notes': slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
            }
        except AttributeError as e:
            print(f"AttributeError: {e}")
            slide_data = {
                'title': slide.shapes.title.text if slide.shapes.title else '',
                'text': '\n'.join(shape.text for shape in slide.shapes if hasattr(shape, 'text')),
                'notes': ''
            }
        if slide_data['title'] == '':
            # slide_data['title'] = "slide contains " + format(len(slide_data['text'].split('\n'))) + " sections."
            first = slide_data['text'].split('\n')[0]
            if not first.isdigit():
                 slide_data['title'] = first
            else:
                 slide_data['title'] = slide_data['text'].split('\n')[1]

        slide_data['text'] = slide_data['text'].replace('\n','\\\n')

        # update module_no if indicated in current slide:
        if extract_module_number(slide_data['text']) != '':
             module_no = extract_module_number(slide_data['text'])
        if module_no == '':
            slide_data['title'] = slide_data['title']
        else:
            slide_data['title'] = module_no + ", " + slide_data['title']
        slide_contents.append(slide_data)

    return slide_contents

def write_to_markdown(slide_contents, output_file):
    """
    Writes the slide content to a Markdown file.

    Args:
        slide_contents (list): A list of dictionaries containing the slide content.
        output_file (str): The path to the output Markdown file.
    """
    with open(output_file, 'w', encoding='utf-8') as file:
        for i, slide in enumerate(slide_contents, start=1):
            file.write(f"#### Slide {i}: {slide['title']}\n\n")
            # file.write(f"### Title: {slide['title']}\n\n")
            # file.write(f"### Text:\n{slide['text']}\n\n")
            file.write(f"{slide['text']}\n\n")
            if slide['notes']:
                file.write(f"### Presenter notes:\n{slide['notes']}\n\n")
            file.write("---\n\n")

slide_contents = extract_slide_content(pptx_file)
write_to_markdown(slide_contents, output_file)
print(f"Slide content written to {output_file}")%
